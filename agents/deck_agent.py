import json
import os
from datetime import datetime
from typing import Any, Dict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from agents.schemas import PitchDeckOutput
from tools.llm_client import call_llm
from tools.structured_output import StructuredOutputError, call_validated_json


class PitchDeckAgent:
    """
    YC-style 12-slide deck generator.
    Produces both JSON slide structure AND a real .pptx file.
    """

    def __init__(self, model=None):
        self.model = model

    def _build_prompt(self, bundle: Dict[str, Any]):
        return f"""
You are a world-class YC-style investor and pitch deck designer.

Using the structured data below, create a **12-slide YC-style pitch deck**.
Do NOT invent financials or metrics. Use only what's provided.

FORMAT: Return ONLY JSON:
{{
  "slides": [
    {{"title": "Problem", "bullets": [], "source_refs": []}},
    {{"title": "Target User", "bullets": [], "source_refs": []}},
    {{"title": "Current Behavior", "bullets": [], "source_refs": []}},
    {{"title": "Solution", "bullets": [], "source_refs": []}},
    {{"title": "Why Now", "bullets": [], "source_refs": []}},
    {{"title": "Market Size", "bullets": [], "source_refs": []}},
    {{"title": "Competition", "bullets": [], "source_refs": []}},
    {{"title": "Unique Advantage", "bullets": [], "source_refs": []}},
    {{"title": "Business Model", "bullets": [], "source_refs": []}},
    {{"title": "Traction", "bullets": [], "source_refs": []}},
    {{"title": "Financial Projection Summary", "bullets": [], "source_refs": []}},
    {{"title": "The Ask (Fundraising)", "bullets": [], "source_refs": []}}
  ]
}}

RULES:
- Return exactly these 12 titles in exactly this order, with at most 6 bullets each.
- Every quantitative claim must be present in the supplied data.
- Put supporting field paths or market URLs in source_refs; do not display them as bullets.
- If evidence for a slide is missing, state the gap plainly instead of inventing content.

========================
STARTUP DATA
========================
{json.dumps(bundle.get("extracted", {}), indent=2)}

========================
MARKET
========================
{json.dumps(bundle.get("market", {}), indent=2)}

========================
FINANCIALS
========================
{json.dumps(bundle.get("financial", {}), indent=2)}

========================
MEMO
========================
{json.dumps(bundle.get("memo", {}), indent=2)}
"""

    def _clean_json(self, text: str) -> str:
        clean = text.replace("```json", "").replace("```", "").strip()
        start = clean.find("{")
        end = clean.rfind("}") + 1
        if start != -1 and end > start:
            return clean[start:end]
        return clean

    # ----------------------------------------------------------------------
    #   PPTX Creation Logic
    # ----------------------------------------------------------------------
    def _create_pptx(self, slides_json: Dict[str, Any], output_path=None):

        # Ensure output folder exists
        os.makedirs("outputs/decks", exist_ok=True)

        # Use timestamp-based filename unless specified
        if output_path is None:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"outputs/decks/pitch_deck_{ts}.pptx"

        prs = Presentation()

        for slide in slides_json["slides"]:
            layout = prs.slide_layouts[1]  # Title + body
            s = prs.slides.add_slide(layout)

            # Title
            s.shapes.title.text = slide["title"]

            # Body
            body = s.placeholders[1].text_frame
            body.clear()

            bullets = slide["bullets"] or ["Evidence not available in source material."]
            for index, bullet in enumerate(bullets):
                # ``clear`` retains one paragraph; reuse it to avoid a blank
                # bullet at the top of every generated slide.
                p = body.paragraphs[0] if index == 0 else body.add_paragraph()
                p.text = str(bullet)
                p.level = 0
                p.font.size = Pt(20)

            if slide.get("source_refs"):
                footer = s.shapes.add_textbox(
                    Inches(0.5),
                    Inches(7.0),
                    Inches(9.0),
                    Inches(0.3),
                ).text_frame
                footer.text = "Sources: " + " · ".join(slide["source_refs"][:3])
                footer.paragraphs[0].font.size = Pt(8)
                footer.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

        prs.save(output_path)
        return output_path

    # ----------------------------------------------------------------------
    #   Main Run Logic
    # ----------------------------------------------------------------------
    def run(self, extracted, market, financial, memo):

        bundle = {
            "extracted": extracted,
            "market": market,
            "financial": financial,
            "memo": memo
        }

        prompt = self._build_prompt(bundle)
        try:
            slides_json = call_validated_json(
                prompt,
                PitchDeckOutput,
                lambda current_prompt: call_llm(current_prompt, model=self.model),
                attempts=2,
            ).model_dump()
        except StructuredOutputError as exc:
            return {
                "error": "Could not validate 12-slide deck JSON",
                "raw": exc.last_response,
                "validation_error": str(exc),
            }

        # Generate PPTX
        pptx_path = self._create_pptx(slides_json)

        return {
            "slides_json": slides_json,
            "pptx_path": pptx_path
        }
