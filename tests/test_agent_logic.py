import json

from pptx import Presentation

import agents.deck_agent as deck_module
import agents.extractor_agent as extractor_module
import agents.market_agent as market_module
import agents.skeptic_agent as skeptic_module
import tools.mcp_client as mcp_client
from agents.deck_agent import PitchDeckAgent
from agents.extractor_agent import ExtractionAgent
from agents.market_agent import MarketAgent
from agents.schemas import SLIDE_TITLES
from agents.skeptic_agent import SkepticAgent


def test_extractor_normalizes_agent_output(monkeypatch):
    response = {
        "name": "Ledgerly",
        "problem": "Manual reconciliation",
        "solution": "Automated reconciliation",
        "target_customer": "Controllers",
        "business_model": "SaaS",
        "pricing": "$2,000 monthly",
        "gtm_strategy": "Direct sales",
        "team": "Finance and ML operators",
        "cost_structure": "Cloud and sales",
        "competition": "Legacy spreadsheets",
        "notable_metrics": [
            "Last month revenue: $42,000",
            {"Monthly active users": "1,200"},
        ],
        "assumptions": "Deck-provided values only",
    }
    monkeypatch.setattr(extractor_module, "call_llm", lambda prompt: json.dumps(response))

    result = ExtractionAgent().extract_from_text("Ledgerly pitch deck")

    assert result["name"] == "Ledgerly"
    assert result["competition"] == ["Legacy spreadsheets"]
    assert result["notable_metrics"]["revenue_last_month"] == "$42,000"
    assert result["notable_metrics"]["mau"] == "1,200"
    assert result["missing_info"] == []


def test_extractor_returns_complete_fallback_for_invalid_json(monkeypatch):
    monkeypatch.setattr(extractor_module, "call_llm", lambda prompt: "not json")

    result = ExtractionAgent().extract_from_text(
        "Nimbus\nBuilds forecasting software.",
        fallback_name="nimbus-deck.pdf",
    )

    assert result["name"] == "Nimbus"
    assert result["team"] == ""
    assert result["competition"] == []
    assert result["notable_metrics"] == {}


def test_extractor_retries_invalid_structured_output(monkeypatch):
    valid = {
        "name": "RetryWorks",
        "problem": "Manual work",
        "solution": "Automation",
    }
    responses = iter(["invalid", json.dumps(valid)])
    monkeypatch.setattr(extractor_module, "call_llm", lambda prompt: next(responses))

    result = ExtractionAgent().extract_from_text("RetryWorks\nAutomation platform")

    assert result["name"] == "RetryWorks"
    assert result["solution"] == "Automation"


def test_market_agent_merges_grounded_sources(monkeypatch):
    llm_output = {
        "market_category": "Finance software",
        "sources": [],
        "competitive_landscape": {},
    }
    web_sources = [
        {"title": "Market study", "url": "https://example.com/market", "snippet": "Large market"},
        {"title": "Growth study", "url": "https://example.com/growth", "snippet": "Growing"},
    ]
    monkeypatch.setenv("DISABLE_WEB_SEARCH", "false")
    monkeypatch.setattr(market_module, "call_llm", lambda prompt: json.dumps(llm_output))
    monkeypatch.setattr(mcp_client, "get_public_comps", lambda ticker: "P/S: 8x")

    result = MarketAgent(use_web_search=True).run(
        {"solution": "Finance automation", "business_model": "SaaS"},
        search_tool=lambda extracted: ("grounded findings", web_sources),
    )

    assert result["web_search_used"] is True
    assert result["public_comparable"]["ticker"] == "CRM"
    assert [source["url"] for source in result["sources"]] == [
        "https://example.com/market",
        "https://example.com/growth",
    ]


def test_skeptic_agent_returns_structured_parse_error(monkeypatch):
    monkeypatch.setattr(skeptic_module, "call_llm", lambda prompt: "invalid response")

    result = SkepticAgent().run({}, {}, {})

    assert result["error"] == "Failed to parse skeptic JSON"
    assert result["raw_response"] == "invalid response"


def test_deck_agent_transforms_llm_json_into_output(monkeypatch, tmp_path):
    slides = {
        "slides": [
            {
                "title": title,
                "bullets": ["Manual diligence is slow"],
                "source_refs": [],
            }
            for title in SLIDE_TITLES
        ]
    }
    output_path = tmp_path / "deck.pptx"
    monkeypatch.setattr(deck_module, "call_llm", lambda prompt, model=None: json.dumps(slides))
    monkeypatch.setattr(
        PitchDeckAgent,
        "_create_pptx",
        lambda self, slides_json: str(output_path),
    )

    result = PitchDeckAgent().run({}, {}, {}, {})

    assert result["slides_json"] == slides
    assert result["pptx_path"] == str(output_path)


def test_deck_renderer_avoids_blank_bullets_and_adds_source_footer(tmp_path):
    slides = {
        "slides": [
            {
                "title": title,
                "bullets": ["Evidence-backed statement"],
                "source_refs": ["extracted.problem"],
            }
            for title in SLIDE_TITLES
        ]
    }
    output_path = tmp_path / "rendered.pptx"

    PitchDeckAgent()._create_pptx(slides, output_path=str(output_path))
    presentation = Presentation(output_path)

    assert len(presentation.slides) == 12
    first_slide_text = [
        shape.text
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text")
    ]
    assert "Evidence-backed statement" in first_slide_text
    assert any(text.startswith("Sources: extracted.problem") for text in first_slide_text)
